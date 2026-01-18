"""
PPTXファイルから構造的なデザイン情報を抽出するユーティリティ。

python-pptxとlxmlを使用してPPTXのテーマ情報（カラースキーム、フォントスキーム、
レイアウト情報）を抽出する。
"""
import logging
from io import BytesIO
from typing import List

from pptx import Presentation
from pptx.opc.constants import RELATIONSHIP_TYPE as RT
from lxml import etree

from src.schemas.design import (
    ColorScheme, FontScheme, 
    SlideLayoutInfo, LayoutPlaceholder, BackgroundInfo,
    LayoutType
)

logger = logging.getLogger(__name__)

# XML名前空間
NSMAP = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'
}


def extract_color_scheme(theme_xml: bytes) -> ColorScheme:
    """テーマXMLからカラースキームを抽出
    
    Args:
        theme_xml: テーマXMLのバイトデータ
        
    Returns:
        ColorScheme: 抽出されたカラースキーム
    """
    root = etree.fromstring(theme_xml)
    
    color_map: dict[str, str] = {}
    scheme_elements = root.xpath('//a:clrScheme/*', namespaces=NSMAP)
    
    for elem in scheme_elements:
        color_name = etree.QName(elem.tag).localname
        
        # srgbClr (RGB直接指定) を探す
        srgb = elem.find('.//a:srgbClr', namespaces=NSMAP)
        if srgb is not None:
            val = srgb.get('val')
            if val:
                color_map[color_name] = f"#{val}"
            continue
        
        # sysClr (システムカラー) を探す
        sys_clr = elem.find('.//a:sysClr', namespaces=NSMAP)
        if sys_clr is not None:
            last_clr = sys_clr.get('lastClr')
            if last_clr:
                color_map[color_name] = f"#{last_clr}"
    
    return ColorScheme(
        dk1=color_map.get('dk1', '#000000'),
        dk2=color_map.get('dk2', '#333333'),
        lt1=color_map.get('lt1', '#FFFFFF'),
        lt2=color_map.get('lt2', '#F5F5F5'),
        accent1=color_map.get('accent1', '#4472C4'),
        accent2=color_map.get('accent2', '#ED7D31'),
        accent3=color_map.get('accent3', '#A5A5A5'),
        accent4=color_map.get('accent4', '#FFC000'),
        accent5=color_map.get('accent5', '#5B9BD5'),
        accent6=color_map.get('accent6', '#70AD47'),
        hlink=color_map.get('hlink', '#0563C1'),
        folHlink=color_map.get('folHlink', '#954F72')
    )


def extract_font_scheme(theme_xml: bytes) -> FontScheme:
    """テーマXMLからフォントスキームを抽出
    
    Args:
        theme_xml: テーマXMLのバイトデータ
        
    Returns:
        FontScheme: 抽出されたフォントスキーム
    """
    root = etree.fromstring(theme_xml)
    
    major_latin = root.xpath('//a:majorFont/a:latin/@typeface', namespaces=NSMAP)
    minor_latin = root.xpath('//a:minorFont/a:latin/@typeface', namespaces=NSMAP)
    major_ea = root.xpath('//a:majorFont/a:ea/@typeface', namespaces=NSMAP)
    minor_ea = root.xpath('//a:minorFont/a:ea/@typeface', namespaces=NSMAP)
    
    return FontScheme(
        major_latin=major_latin[0] if major_latin else "Calibri Light",
        minor_latin=minor_latin[0] if minor_latin else "Calibri",
        major_east_asian=major_ea[0] if major_ea else None,
        minor_east_asian=minor_ea[0] if minor_ea else None
    )


def _infer_layout_type(layout_name: str) -> LayoutType:
    """レイアウト名からレイアウトタイプを推定
    
    Args:
        layout_name: レイアウト名
        
    Returns:
        LayoutType: 推定されたレイアウトタイプ
    """
    name_lower = layout_name.lower()
    
    if "title slide" in name_lower or "タイトル スライド" in name_lower:
        return "title_slide"
    elif "section" in name_lower or "セクション" in name_lower:
        return "section_header"
    elif "comparison" in name_lower or "比較" in name_lower:
        return "comparison"
    elif "two content" in name_lower or "2 つのコンテンツ" in name_lower:
        return "two_content"
    elif "picture" in name_lower or "図" in name_lower:
        return "picture_with_caption"
    elif "blank" in name_lower or "白紙" in name_lower:
        return "blank"
    elif "content" in name_lower or "コンテンツ" in name_lower:
        return "title_and_content"
    else:
        return "other"


def extract_layout_info(layout, slide_width: int, slide_height: int) -> SlideLayoutInfo:
    """スライドレイアウトからレイアウト情報を抽出
    
    Args:
        layout: python-pptxのSlideLayoutオブジェクト
        slide_width: スライドの幅（EMU単位）
        slide_height: スライドの高さ（EMU単位）
        
    Returns:
        SlideLayoutInfo: 抽出されたレイアウト情報
    """
    placeholders: List[LayoutPlaceholder] = []
    
    for ph in layout.placeholders:
        left_pct = (ph.left / slide_width * 100) if ph.left else 0
        top_pct = (ph.top / slide_height * 100) if ph.top else 0
        width_pct = (ph.width / slide_width * 100) if ph.width else 0
        height_pct = (ph.height / slide_height * 100) if ph.height else 0
        
        # プレースホルダータイプのマッピング
        ph_type_map = {
            1: "title", 2: "body", 3: "subtitle", 
            15: "title", 18: "picture", 19: "chart",
            20: "table", 11: "footer", 13: "slide_number", 12: "date"
        }
        ph_type = ph_type_map.get(ph.placeholder_format.type, "body")
        
        placeholders.append(LayoutPlaceholder(
            type=ph_type,
            left_percent=round(left_pct, 2),
            top_percent=round(top_pct, 2),
            width_percent=round(width_pct, 2),
            height_percent=round(height_pct, 2)
        ))
    
    layout_type = _infer_layout_type(layout.name)
    
    return SlideLayoutInfo(
        name=layout.name,
        layout_type=layout_type,
        placeholders=placeholders,
        template_image_url=None
    )


def extract_background_info(slide_master) -> BackgroundInfo:
    """スライドマスターから背景情報を抽出
    
    Args:
        slide_master: python-pptxのSlideMasterオブジェクト
        
    Returns:
        BackgroundInfo: 抽出された背景情報
    """
    try:
        background = slide_master.background
        fill = background.fill
        
        if fill.type is None:
            return BackgroundInfo(fill_type="none")
        
        fill_type_map = {1: "solid", 3: "gradient", 6: "picture", 2: "pattern"}
        fill_type = fill_type_map.get(fill.type, "none")
        
        solid_color = None
        if fill.type == 1:
            try:
                rgb = fill.fore_color.rgb
                solid_color = f"#{rgb}"
            except Exception:
                pass
        
        return BackgroundInfo(fill_type=fill_type, solid_color=solid_color)
    except Exception as e:
        logger.warning(f"Background extraction failed: {e}")
        return BackgroundInfo(fill_type="none")


def extract_design_context_structure(pptx_bytes: bytes, filename: str = "template.pptx") -> tuple[
    ColorScheme, FontScheme, List[SlideLayoutInfo], BackgroundInfo, int, int
]:
    """PPTXバイトデータから構造的なデザイン情報を抽出（画像レンダリングなし）
    
    Args:
        pptx_bytes: PPTXファイルのバイトデータ
        filename: ファイル名（ログ用）
        
    Returns:
        tuple: (color_scheme, font_scheme, layouts, background, slide_master_count, layout_count)
    """
    logger.info(f"Extracting design context from: {filename}")
    prs = Presentation(BytesIO(pptx_bytes))
    
    # テーマXMLを取得
    slide_master = prs.slide_masters[0]
    theme_part = slide_master.part.part_related_by(RT.THEME)
    theme_xml = theme_part.blob
    
    # 各要素を抽出
    color_scheme = extract_color_scheme(theme_xml)
    font_scheme = extract_font_scheme(theme_xml)
    background = extract_background_info(slide_master)
    
    # スライドサイズを取得
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 全レイアウトを抽出
    layouts: List[SlideLayoutInfo] = []
    for layout in prs.slide_layouts:
        layout_info = extract_layout_info(layout, slide_width, slide_height)
        layouts.append(layout_info)
    
    logger.info(
        f"Extraction complete: {len(layouts)} layouts, "
        f"accent1={color_scheme.accent1}, font={font_scheme.major_latin}"
    )
    
    return (
        color_scheme, 
        font_scheme, 
        layouts, 
        background, 
        len(prs.slide_masters), 
        len(prs.slide_layouts)
    )
